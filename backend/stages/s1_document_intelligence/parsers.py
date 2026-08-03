"""Format parsers: PDF, DOCX, PPTX, plain text.

Each returns ordered :class:`Block` objects. Nothing here interprets meaning —
that is stage 2 onward. This layer's whole job is to lose as little structure as
possible on the way in, because no later stage can recover what parsing discards.

Kept as one module rather than a package: the four parsers are ~50 lines each and
share the block-emitting helpers, so splitting them would add navigation cost
without adding separation.

DOCX and PPTX carry explicit styles, so their headings and tables are read
directly and are reliable. PDF has no markup, so structure is inferred from
typography — see ``structure.py`` for the heuristics and their limits.

Every parser here runs on attacker-supplied bytes, so all four take the *same*
explicit limits and none of them accepts ``**kwargs``. A parser that silently
swallowed a limit it was handed is the bug that signature shape exists to make
impossible.
"""

from __future__ import annotations

import io
import zipfile
from collections.abc import Iterator
from typing import Any

from contracts.document import Block, TableData
from stages.s1_document_intelligence.equations import is_probable_equation, to_latex
from stages.s1_document_intelligence.errors import (
    DocumentTooLarge,
    IngestionError,
    ParseFailure,
    TooManyPages,
)
from stages.s1_document_intelligence.structure import (
    TextSpan,
    body_font_size,
    infer_heading_level,
)

__all__ = [
    "DEFAULT_MAX_ARCHIVE_BYTES",
    "DEFAULT_MAX_ARCHIVE_MEMBERS",
    "DEFAULT_MAX_ARCHIVE_RATIO",
    "DEFAULT_MAX_BLOCKS",
    "DEFAULT_MAX_TEXT_CHARS",
    "OOXML_MIMES",
    "inspect_ooxml_archive",
    "parse_docx",
    "parse_pdf",
    "parse_pptx",
    "parse_text",
]

_LIST_MARKERS = ("- ", "* ", "• ", "– ")

#: The two OOXML container types this pipeline accepts. Both are ZIP archives,
#: which is what makes the archive inspection below necessary.
OOXML_MIMES = frozenset(
    {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
)

#: Extraction ceilings. Calibrated against the benchmark input named in FAQ.md —
#: an NCERT Class 11 Physics chapter (44 pages) yields ~2.7k blocks and ~98k
#: characters, so these sit two to three orders of magnitude above a legitimate
#: chapter and only a document engineered to exhaust memory can reach them.
DEFAULT_MAX_BLOCKS = 200_000
DEFAULT_MAX_TEXT_CHARS = 20_000_000

#: Archive-bomb thresholds. See :func:`inspect_ooxml_archive`.
DEFAULT_MAX_ARCHIVE_BYTES = 200 * 1024 * 1024
DEFAULT_MAX_ARCHIVE_RATIO = 200
DEFAULT_MAX_ARCHIVE_MEMBERS = 2000

#: A member smaller than this cannot exhaust anything however well it compressed,
#: so the ratio rule ignores it. Without this floor a small run of repeated
#: whitespace — which real Word documents contain — would trip the check.
_RATIO_FLOOR_BYTES = 1024 * 1024


# ─────────────────────────────────────────────────── archive bomb inspection ───


def inspect_ooxml_archive(
    payload: bytes,
    *,
    max_total_bytes: int = DEFAULT_MAX_ARCHIVE_BYTES,
    max_ratio: int = DEFAULT_MAX_ARCHIVE_RATIO,
    max_members: int = DEFAULT_MAX_ARCHIVE_MEMBERS,
) -> None:
    """Reject a decompression bomb before any parser opens the container.

    DOCX and PPTX are ZIP archives. ``python-docx`` and ``python-pptx`` hand the
    whole of ``document.xml`` to lxml as one string, so a 2 MB upload declaring a
    3 GB member is a 3 GB allocation in a worker — the upload size limit never
    sees it, because the upload really is 2 MB.

    Only the central directory is read here: no member is decompressed, so the
    check costs microseconds and cannot itself be turned into the attack. Reading
    *declared* sizes is sound rather than merely cheap, because :mod:`zipfile`
    truncates every read at the declared ``file_size`` (``ZipExtFile._read1``),
    so a lying directory yields a short read and a CRC error, not a bomb.

    Three independent rules, because a bomb only has to fail one:

    * **total declared size** — the aggregate a parser could be asked to hold;
    * **per-member ratio** — one hostile member hiding inside a normal archive;
    * **member count** — thousands of small parts, each cheap, ruinous in bulk.

    Raises :class:`DocumentTooLarge` (HTTP 413) on any of the three, and
    :class:`ParseFailure` when the container is not a readable archive.
    """
    try:
        with zipfile.ZipFile(io.BytesIO(payload)) as archive:
            infos = archive.infolist()
    except zipfile.BadZipFile as exc:
        raise ParseFailure(f"Not a readable OOXML container: {exc}") from exc

    members = [info for info in infos if not info.is_dir()]
    if len(members) > max_members:
        raise DocumentTooLarge(
            f"Archive declares {len(members)} members; limit is {max_members}.",
            reason="archive_member_count",
            members=len(members),
            limit_members=max_members,
        )

    total = 0
    for info in members:
        total += info.file_size
        if total > max_total_bytes:
            raise DocumentTooLarge(
                f"Archive expands to at least {total} bytes; limit is {max_total_bytes}.",
                reason="archive_uncompressed_size",
                uncompressed_bytes=total,
                limit_bytes=max_total_bytes,
            )
        if info.file_size < _RATIO_FLOOR_BYTES:
            continue
        # A stored member has a ratio of 1 and is already caught by the total;
        # only a member that actually expands is interesting here.
        ratio = (
            float(info.file_size)
            if info.compress_size <= 0
            else info.file_size / info.compress_size
        )
        if ratio > max_ratio:
            raise DocumentTooLarge(
                f"Archive member {info.filename!r} expands {ratio:.0f}:1; limit is {max_ratio}:1.",
                reason="archive_compression_ratio",
                member=info.filename,
                ratio=round(ratio, 1),
                limit_ratio=max_ratio,
                uncompressed_bytes=info.file_size,
            )


class _BlockBuilder:
    """Assigns ids and character offsets so blocks stay addressable.

    Also the single choke point for extraction volume. Every block of every
    format is created here, so capping here caps all four parsers uniformly —
    the alternative, a per-parser check, is four places to forget.
    """

    def __init__(
        self,
        *,
        max_blocks: int = DEFAULT_MAX_BLOCKS,
        max_chars: int = DEFAULT_MAX_TEXT_CHARS,
    ) -> None:
        self.blocks: list[Block] = []
        self._cursor = 0
        self._max_blocks = max_blocks
        self._max_chars = max_chars

    def add(
        self,
        block_type: str,
        text: str,
        *,
        page: int | None = None,
        level: int | None = None,
        latex: str | None = None,
        table: TableData | None = None,
    ) -> None:
        text = text.strip()
        if not text and table is None:
            return

        # Checked before the block is built, so the ceiling bounds what is held
        # rather than what has already been held. Truncating silently would be
        # worse than rejecting: a package built from half a document still looks
        # complete, and nothing downstream can tell that it is not.
        if len(self.blocks) >= self._max_blocks:
            raise DocumentTooLarge(
                f"Document yields more than {self._max_blocks} blocks.",
                reason="block_ceiling",
                limit_blocks=self._max_blocks,
            )
        if self._cursor + len(text) > self._max_chars:
            raise DocumentTooLarge(
                f"Document yields more than {self._max_chars} characters of text.",
                reason="character_ceiling",
                limit_chars=self._max_chars,
            )

        start = self._cursor
        self._cursor += len(text) + 1
        self.blocks.append(
            Block(
                block_id=f"b_{len(self.blocks):04d}",
                type=block_type,  # type: ignore[arg-type]
                text=text or "[table]",
                page=page,
                level=level,
                latex=latex,
                table=table,
                char_start=start,
                char_end=self._cursor - 1,
            )
        )


def _classify_line(builder: _BlockBuilder, span: TextSpan, body_size: float | None) -> None:
    """Route one line to the right block type. Order matters: equations first."""
    text = span.text.strip()
    if not text:
        return

    if is_probable_equation(text):
        builder.add("equation", text, page=span.page, latex=to_latex(text))
        return

    level = infer_heading_level(span, body_size)
    if level is not None:
        builder.add("heading", text, page=span.page, level=level)
        return

    if text.startswith(_LIST_MARKERS) or (
        len(text) > 2 and text[0].isdigit() and text[1:3] in (". ", ") ")
    ):
        builder.add("list", text, page=span.page)
        return

    builder.add("paragraph", text, page=span.page)


# ───────────────────────────────────────────────────────────────────── PDF ───


def _pdf_spans(
    page: Any, page_no: int, table_boxes: list[tuple[float, float, float, float]]
) -> Iterator[TextSpan]:
    """Reconstruct lines with their dominant font size.

    pdfplumber yields positioned characters; lines are rebuilt by grouping on the
    vertical midpoint, because PDFs have no concept of a line.

    Words inside a table's bounding box are excluded. Without this, a bold table
    header row satisfies every heading cue — short, unpunctuated, heavier than
    body text — and gets promoted into the outline, corrupting the section paths
    of everything that follows it.
    """

    def inside_table(word: dict[str, Any]) -> bool:
        cx = (word["x0"] + word["x1"]) / 2
        cy = (word["top"] + word["bottom"]) / 2
        return any(x0 <= cx <= x1 and top <= cy <= bottom for x0, top, x1, bottom in table_boxes)

    words = [
        w
        for w in (page.extract_words(extra_attrs=["size", "fontname"]) or [])
        if not inside_table(w)
    ]
    rows: dict[int, list[dict[str, Any]]] = {}
    for word in words:
        key = round((word["top"] + word["bottom"]) / 2)
        rows.setdefault(key, []).append(word)

    for key in sorted(rows):
        row = sorted(rows[key], key=lambda w: w["x0"])
        text = " ".join(w["text"] for w in row)
        sizes = [w.get("size") for w in row if w.get("size")]
        fonts = [str(w.get("fontname", "")) for w in row]
        yield TextSpan(
            text=text,
            page=page_no,
            font_size=max(sizes) if sizes else None,
            bold=any("bold" in f.lower() for f in fonts),
        )


def parse_pdf(
    payload: bytes,
    *,
    max_pages: int,
    max_blocks: int = DEFAULT_MAX_BLOCKS,
    max_chars: int = DEFAULT_MAX_TEXT_CHARS,
) -> list[Block]:
    import pdfplumber

    builder = _BlockBuilder(max_blocks=max_blocks, max_chars=max_chars)
    try:
        with pdfplumber.open(io.BytesIO(payload)) as pdf:
            if len(pdf.pages) > max_pages:
                raise TooManyPages(
                    f"Document has {len(pdf.pages)} pages; limit is {max_pages}.",
                    page_count=len(pdf.pages),
                    limit=max_pages,
                )

            # Two passes: collect spans to establish the body font size, then
            # classify. Heading inference is relative, so it needs the whole
            # document before it can judge any single line.
            pages: list[tuple[int, list[TextSpan], list[list[list[str]]]]] = []
            all_spans: list[TextSpan] = []
            for index, page in enumerate(pdf.pages, start=1):
                found = page.find_tables() or []
                boxes = [t.bbox for t in found]
                spans = list(_pdf_spans(page, index, boxes))
                all_spans.extend(spans)
                pages.append((index, spans, [t.extract() for t in found]))

            body_size = body_font_size(all_spans)

            for index, spans, tables in pages:
                for span in spans:
                    _classify_line(builder, span, body_size)
                for raw in tables:
                    _add_table(builder, raw, page=index)
    except IngestionError:
        # A limit was hit. It is already the right rejection carrying the right
        # code and details; re-wrapping it as a parse failure would lose both.
        raise
    except Exception as exc:
        raise ParseFailure(f"Could not parse PDF: {exc}") from exc

    return builder.blocks


def _add_table(builder: _BlockBuilder, raw: list[list[Any]], *, page: int | None) -> None:
    """Emit a table as structure. Ragged rows are padded, never dropped."""
    rows = [[("" if cell is None else str(cell).strip()) for cell in row] for row in raw if row]
    if not rows:
        return
    headers, body = rows[0], rows[1:]
    width = len(headers)
    normalised = [(r + [""] * width)[:width] for r in body]
    preview = " | ".join(headers)
    builder.add(
        "table",
        preview,
        page=page,
        table=TableData(headers=headers, rows=normalised),
    )


# ──────────────────────────────────────────────────────────────────── DOCX ───


def _docx_declared_pages(document: Any) -> int | None:
    """Pages a DOCX *declares*, or ``None`` when it declares none.

    A DOCX has no page count until it is laid out, so this reads the only two
    things the file actually records: explicit page breaks the author inserted,
    and the ``lastRenderedPageBreak`` markers Word writes on save. Returning
    ``None`` when neither is present is deliberate — an estimate derived from
    paragraph counts would reject real documents, and a limit that fires on a
    guess is worse than no limit at all. Volume is still bounded, by the
    builder's block and character ceilings.
    """
    try:
        from docx.oxml.ns import qn  # type: ignore[import-untyped]

        body = document.element.body
        rendered = len(body.findall(f".//{qn('w:lastRenderedPageBreak')}"))
        explicit = sum(1 for br in body.iter(qn("w:br")) if br.get(qn("w:type")) == "page")
    except Exception:
        return None
    breaks = max(rendered, explicit)
    return breaks + 1 if breaks else None


def parse_docx(
    payload: bytes,
    *,
    max_pages: int = 300,
    max_blocks: int = DEFAULT_MAX_BLOCKS,
    max_chars: int = DEFAULT_MAX_TEXT_CHARS,
) -> list[Block]:
    from docx import Document  # type: ignore[import-untyped]

    builder = _BlockBuilder(max_blocks=max_blocks, max_chars=max_chars)
    try:
        document = Document(io.BytesIO(payload))

        declared_pages = _docx_declared_pages(document)
        if declared_pages is not None and declared_pages > max_pages:
            raise TooManyPages(
                f"Document declares {declared_pages} pages; limit is {max_pages}.",
                page_count=declared_pages,
                limit=max_pages,
            )

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style = (paragraph.style.name or "").lower()

            if style.startswith("heading"):
                digits = "".join(ch for ch in style if ch.isdigit())
                builder.add("heading", text, level=int(digits) if digits else 1)
            elif style.startswith("title"):
                builder.add("heading", text, level=1)
            elif is_probable_equation(text):
                builder.add("equation", text, latex=to_latex(text))
            elif "list" in style:
                builder.add("list", text)
            else:
                builder.add("paragraph", text)

        for table in document.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            _add_table(builder, rows, page=None)
    except IngestionError:
        raise
    except Exception as exc:
        raise ParseFailure(f"Could not parse DOCX: {exc}") from exc

    return builder.blocks


# ──────────────────────────────────────────────────────────────────── PPTX ───


def parse_pptx(
    payload: bytes,
    *,
    max_pages: int = 300,
    max_blocks: int = DEFAULT_MAX_BLOCKS,
    max_chars: int = DEFAULT_MAX_TEXT_CHARS,
) -> list[Block]:
    from pptx import Presentation  # type: ignore[import-untyped]

    builder = _BlockBuilder(max_blocks=max_blocks, max_chars=max_chars)
    try:
        presentation = Presentation(io.BytesIO(payload))
        slides = list(presentation.slides)
        if len(slides) > max_pages:
            raise TooManyPages(
                f"Deck has {len(slides)} slides; limit is {max_pages}.",
                page_count=len(slides),
                limit=max_pages,
            )

        for index, slide in enumerate(slides, start=1):
            # The title placeholder is the slide's heading — explicit, so trusted.
            title = getattr(slide.shapes, "title", None)
            if title is not None and title.has_text_frame and title.text.strip():
                builder.add("heading", title.text.strip(), page=index, level=1)

            for shape in slide.shapes:
                if shape is title:
                    continue
                if getattr(shape, "has_table", False):
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    _add_table(builder, rows, page=index)
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if not text:
                        continue
                    if is_probable_equation(text):
                        builder.add("equation", text, page=index, latex=to_latex(text))
                    elif paragraph.level and paragraph.level > 0:
                        builder.add("list", text, page=index)
                    else:
                        builder.add("paragraph", text, page=index)

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    builder.add("paragraph", notes, page=index)
    except IngestionError:
        raise
    except Exception as exc:
        raise ParseFailure(f"Could not parse PPTX: {exc}") from exc

    return builder.blocks


# ─────────────────────────────────────────────────────────────── plain text ───


def parse_text(
    payload: bytes,
    *,
    max_pages: int = 300,
    max_blocks: int = DEFAULT_MAX_BLOCKS,
    max_chars: int = DEFAULT_MAX_TEXT_CHARS,
) -> list[Block]:
    """Plain text and Markdown.

    Markdown ``#`` prefixes are honoured as explicit heading levels; otherwise the
    same shape heuristics apply, with no font size to work from.

    ``max_pages`` is accepted and unused on purpose: plain text has no pages, and
    a uniform signature across the four parsers is what stops a limit being
    dropped on the floor by a ``**kwargs``. Volume is bounded by the builder.
    """
    try:
        content = payload.decode("utf-8")
    except UnicodeDecodeError:
        try:
            content = payload.decode("latin-1")
        except Exception as exc:
            raise ParseFailure("Could not decode text document.") from exc

    builder = _BlockBuilder(max_blocks=max_blocks, max_chars=max_chars)
    for raw_line in content.splitlines():
        line = raw_line.rstrip()
        if not line.strip():
            continue
        if line.lstrip().startswith("#"):
            marker = line.lstrip()
            level = len(marker) - len(marker.lstrip("#"))
            builder.add("heading", marker.lstrip("#").strip(), level=min(level, 6))
            continue
        _classify_line(builder, TextSpan(text=line), body_size=None)

    return builder.blocks
